"""
LangChain Multi-Tools Module for Educational Multi-Agent System

Provides specialized tools for:
- Knowledge retrieval using Elasticsearch RAG
- Quiz generation using text-generation models
- Report generation using text-generation models (NEW)
- PDF document creation
- PPTX presentation generation (NEW)
- Email delivery with attachments

Each tool is decorated with @tool for seamless LangChain integration.
Models are initialized once at module load time.
All responses are JSON-formatted for consistent parsing by the main agent.
"""

import os
import re
import json
import logging
import smtplib
from pathlib import Path
import torch
from email.mime.multipart import MIMEMultipart
from email.mime.text import MIMEText
from email.mime.base import MIMEBase
from email import encoders
from langchain.tools import tool
from fpdf import FPDF
from transformers import AutoTokenizer, AutoModelForCausalLM, BitsAndBytesConfig

# NEW Import: Added for presentation generation
try:
    from pptx import Presentation
    from pptx.util import Inches, Pt
except ImportError:
    print("WARNING: `python-pptx` not installed. Presentation generator will fail.")
    print("Please run: pip install python-pptx")

# ======================== LOGGING CONFIGURATION ========================

def _setup_logger():
    """Configure logger for tools module."""
    log_dir = Path("logs")
    log_dir.mkdir(exist_ok=True)
    
    logger = logging.getLogger("multitools")
    if not logger.handlers:
        logger.setLevel(logging.INFO)
        handler = logging.FileHandler(log_dir / "tools.log", encoding="utf-8")
        formatter = logging.Formatter(
            "%(asctime)s - %(name)s - %(levelname)s - %(message)s"
        )
        handler.setFormatter(formatter)
        logger.addHandler(handler)
    
    return logger

logger = _setup_logger()

# ======================== GLOBAL INITIALIZATION ========================

# RAG Retriever - initialize once
_retriever = None

def _init_retriever():
    """Initialize RAG retriever from retrieval.py module."""
    global _retriever
    if _retriever is None:
        try:
            from retrieval import Retriever
            _retriever = Retriever(log_dir=Path("logs"))
        except Exception as e:
            logger.error(f"Failed to initialize retriever: {e}", exc_info=True)
            _retriever = None
    return _retriever

# Text Generation Model - initialize once
_text_gen_model = None
_text_gen_tokenizer = None
_text_gen_model_name = None
_text_gen_is_qwen = False

def _init_text_gen_model():
    """Initialize text generation model."""
    global _text_gen_model, _text_gen_tokenizer, _text_gen_model_name, _text_gen_is_qwen
    if _text_gen_model is None:
        try:
            model_name = os.getenv("TEXT_GEN_MODEL", "Qwen/Qwen3-1.7B")
            _text_gen_model_name = model_name
            _text_gen_is_qwen = "qwen" in model_name.lower()

            _text_gen_tokenizer = AutoTokenizer.from_pretrained(model_name)

            has_bnb = True
            try:
                import bitsandbytes as bnb
            except ImportError:
                has_bnb = False

            if _text_gen_is_qwen and has_bnb:
                logger.info("Loading text generation model (Qwen) in 8-bit mode.")
                quant_config = BitsAndBytesConfig(
                    load_in_8bit=True,
                    llm_int8_threshold=6.0,
                    llm_int8_enable_fp32_cpu_offload=False,
                )
                _text_gen_model = AutoModelForCausalLM.from_pretrained(
                    model_name,
                    quantization_config=quant_config,
                    device_map="auto",
                )
            else:
                logger.info("Loading full-precision model (CPU or GPU).")
                _text_gen_model = AutoModelForCausalLM.from_pretrained(
                    model_name,
                    torch_dtype=torch.float16 if torch.cuda.is_available() else torch.float32,
                    device_map="auto" if torch.cuda.is_available() else None,
                )
            _text_gen_model.eval()
        except Exception as e:
            logger.error(f"Failed to initialize text generation model: {e}", exc_info=True)
            _text_gen_model = None
            _text_gen_tokenizer = None
    return _text_gen_model, _text_gen_tokenizer

def _tg_move_to(device: str = "cpu"):
    try:
        if _text_gen_model is not None:
            _text_gen_model.to(device)
    except Exception as e:
        logger.warning(f"Failed to move text gen model to {device}: {e}")

def _tg_offload_gpu():
    """Move text generation model back to CPU and clear CUDA cache."""
    _tg_move_to("cpu")
    try:
        if torch.cuda.is_available():
            torch.cuda.empty_cache()
    except Exception:
        pass

# ======================== TOOLS ========================

@tool
def knowledge_retrieval(query: str) -> list[str]:
    """
    Retrieve chemistry knowledge from indexed database using RAG.
    ...
    """
    try:
        if not query or len(query.strip()) < 3:
            error_msg = "Query must be at least 3 characters long."
            logger.warning(f"knowledge_retrieval: {error_msg}")
            return [f"Error: {error_msg}"]
        
        retriever = _init_retriever()
        if retriever is None:
            error_msg = "RAG retriever not initialized. Check Elasticsearch connection."
            logger.error(f"knowledge_retrieval: {error_msg}")
            return ["Error: Knowledge retrieval unavailable. Check logs for details."]
        
        results = retriever.search(query, top_k=20, rerank_top_n=5)
        
        if not results:
            logger.warning(f"knowledge_retrieval: No results found for query: {query}")
            return ["No relevant information found in knowledge base."]
        
        context_chunks = []
        for chunk in results:
            if "_source" in chunk and "chunk_text" in chunk["_source"]:
                context_chunks.append(chunk["_source"]["chunk_text"])
        
        if not context_chunks:
            error_msg = "Retrieved chunks have unexpected format."
            logger.error(f"knowledge_retrieval: {error_msg}")
            return ["Error: Retrieved content format invalid. Check logs for details."]
        
        logger.info(f"knowledge_retrieval: Successfully retrieved {len(context_chunks)} chunks for query: {query}")
        
        # Return a single joined string of knowledge, easier for LLMs
        joined_context = "\n\n---\n\n".join(context_chunks)
        return {"status": "success", "result": joined_context}
        
    except Exception as e:
        logger.error(f"knowledge_retrieval exception: {e}", exc_info=True)
        return {"status": "error", "message": "Error: Knowledge retrieval failed. Check logs for details."}


@tool
def quiz_generator(context: str, constraints: str) -> str:
    """
    Generate educational quiz content from knowledge context.
    ...
    """
    try:
        if not context or len(context.strip()) < 50:
            error_msg = "Context must be at least 50 characters long."
            logger.warning(f"quiz_generator: {error_msg}")
            return json.dumps({"status": "error", "message": error_msg})
        
        # ... (rest of validation) ...
        
        model, tokenizer = _init_text_gen_model()
        if model is None or tokenizer is None:
            # ... (error handling) ...
            pass
        
        # Build prompt for quiz generation
        prompt = f"""Generate an educational quiz.
Requirements: {constraints}
Knowledge Base:
{context[:2000]}
...
Begin quiz:"""
        
        if _text_gen_is_qwen:
            messages = [
                {"role": "system", "content": ("You generate high-quality, well-formatted educational quizzes.")},
                {"role": "user", "content": (
                    f"Constraints: {constraints}\n\n"
                    f"Context:\n{context[:4000]}\n\n"
                    "Generate a numbered multiple-choice quiz (A-D options) with answers labeled 'Answer: [letter]'."
                )}
            ]
            prompt = _text_gen_tokenizer.apply_chat_template(
                messages, tokenize=False, add_generation_prompt=True, enable_thinking=True
            )
        
        inputs = _text_gen_tokenizer([prompt], return_tensors="pt", truncation=True, max_length=2048)

        # ... (offload/move to device logic) ...
        try:
            if _retriever is not None: _retriever.offload_models_to_cpu()
        except Exception: pass
        compute_device = "cuda" if torch.cuda.is_available() else "cpu"
        _tg_move_to(compute_device)
        try: model_device = next(model.parameters()).device
        except Exception: model_device = torch.device(compute_device)
        inputs = {k: v.to(model_device) for k, v in inputs.items()}

        with torch.no_grad():
            # ... (generation logic) ...
            if _text_gen_is_qwen:
                generated_ids = model.generate(
                    **inputs, max_new_tokens=1024, do_sample=True, temperature=0.7, top_p=0.95
                )
                output_ids = generated_ids[0][len(inputs['input_ids'][0]):].tolist()
                try: index = len(output_ids) - output_ids[::-1].index(151668)
                except ValueError: index = 0
                quiz_content = _text_gen_tokenizer.decode(output_ids[index:], skip_special_tokens=True).strip()
            else:
                outputs = model.generate(
                    **inputs, max_new_tokens=1024, num_beams=4, early_stopping=True, temperature=0.7, do_sample=True
                )
                quiz_content = _text_gen_tokenizer.decode(outputs[0][len(inputs['input_ids'][0]):], skip_special_tokens=True)
        
        # ... (validation of output) ...
        
        _tg_offload_gpu()

        return json.dumps({
            "status": "success",
            "result": quiz_content,
            "length": len(quiz_content)
        })
        
    except Exception as e:
        logger.error(f"quiz_generator exception: {e}", exc_info=True)
        return json.dumps({"status": "error", "message": "Error: Quiz generation failed."})

# --- NEW TOOL ---
@tool
def report_generator(context: str, topic: str) -> str:
    """
    Generates a written report/summary from a knowledge context.
    
    Uses a text-generation model to synthesize the provided context
    into a coherent, multi-paragraph report on the given topic.
    
    Args:
        context: Knowledge content to generate the report from.
        topic: The central topic of the report (e.g., "The French Revolution").
        
    Returns:
        JSON string with status and the generated report, or error message.
    """
    try:
        if not context or len(context.strip()) < 50:
            error_msg = "Context must be at least 50 characters long."
            logger.warning(f"report_generator: {error_msg}")
            return json.dumps({"status": "error", "message": error_msg})
        
        model, tokenizer = _init_text_gen_model()
        if model is None or tokenizer is None:
            error_msg = "Text generation model not initialized."
            logger.error(f"report_generator: {error_msg}")
            return json.dumps({"status": "error", "message": "Error: Report generation unavailable."})
        
        # Build prompt for report generation
        prompt = f"""Generate a detailed, well-structured report.
Topic: {topic}
Knowledge Base:
{context[:4000]}

Write a comprehensive report based *only* on the knowledge base provided.
The report should have an introduction, several body paragraphs, and a conclusion.
Format it clearly with paragraphs.

Begin report:"""
        
        if _text_gen_is_qwen:
            messages = [
                {"role": "system", "content": ("You are an expert report writer. You synthesize context into clear, professional, multi-paragraph reports.")},
                {"role": "user", "content": (
                    f"Topic: {topic}\n\n"
                    f"Context:\n{context[:4000]}\n\n"
                    "Write a comprehensive report based *only* on the provided context. Use an introduction, body, and conclusion structure."
                )}
            ]
            prompt = _text_gen_tokenizer.apply_chat_template(
                messages, tokenize=False, add_generation_prompt=True, enable_thinking=True
            )
        
        inputs = _text_gen_tokenizer([prompt], return_tensors="pt", truncation=True, max_length=4096)

        # Offload/move to device logic
        try:
            if _retriever is not None: _retriever.offload_models_to_cpu()
        except Exception: pass
        compute_device = "cuda" if torch.cuda.is_available() else "cpu"
        _tg_move_to(compute_device)
        try: model_device = next(model.parameters()).device
        except Exception: model_device = torch.device(compute_device)
        inputs = {k: v.to(model_device) for k, v in inputs.items()}

        with torch.no_grad():
            if _text_gen_is_qwen:
                generated_ids = model.generate(
                    **inputs, max_new_tokens=2048, do_sample=True, temperature=0.7, top_p=0.95
                )
                output_ids = generated_ids[0][len(inputs['input_ids'][0]):].tolist()
                try: index = len(output_ids) - output_ids[::-1].index(151668)
                except ValueError: index = 0
                report_content = _text_gen_tokenizer.decode(output_ids[index:], skip_special_tokens=True).strip()
            else:
                outputs = model.generate(
                    **inputs, max_new_tokens=2048, num_beams=4, early_stopping=True, temperature=0.7, do_sample=True
                )
                report_content = _text_gen_tokenizer.decode(outputs[0][len(inputs['input_ids'][0]):], skip_special_tokens=True)
                
        if not report_content or len(report_content) < 50:
            error_msg = "Report generation produced insufficient content."
            logger.warning(f"report_generator: {error_msg}")
            return json.dumps({"status": "error", "message": "Error: Generated report is too short."})
        
        logger.info(f"report_generator: Successfully generated report ({len(report_content)} characters) for topic: {topic}")

        _tg_offload_gpu()

        return json.dumps({
            "status": "success",
            "result": report_content,
            "length": len(report_content)
        })
        
    except Exception as e:
        logger.error(f"report_generator exception: {e}", exc_info=True)
        return json.dumps({"status": "error", "message": "Error: Report generation failed."})

@tool
def pdf_generator(content: str, filename: str, title: str):
    """Generate a clean, printable PDF for quiz or report content."""
    try:
        # content = inputs.get("content", "") # No longer need this
        # filename = inputs.get("filename", "output.pdf") # No longer need this
        # title = inputs.get("title", "Generated Document") # No longer need this

        # Ensure output directory exists
        script_dir = os.path.dirname(os.path.abspath(__file__))
        output_dir = os.path.abspath(os.path.join(script_dir, "..", "generated_documents"))
        os.makedirs(output_dir, exist_ok=True)
        filepath = os.path.join(output_dir, filename)

        # Clean up
        content = content.replace("**", "")
        content = content.replace("\\n", "\n").replace("\\t", " ")
        content = content.replace("\r", "")
        content = "\n".join([line.strip() for line in content.splitlines() if line.strip()])

        def safe_chunk(line, limit=100):
            # ... (safe chunk logic) ...
            pass

        pdf = FPDF()
        pdf.add_page()
        pdf.set_auto_page_break(auto=True, margin=15)
        pdf.set_font("Helvetica", "", 12)

        # Header
        pdf.set_font("Helvetica", "B", 16)
        pdf.cell(0, 10, title, align="C", ln=True)
        pdf.ln(8)
        pdf.set_font("Helvetica", "", 12)

        # Render each line safely
        for raw_line in content.split("\n"):
            # ... (rendering logic) ...
            try:
                pdf.multi_cell(0, 8, raw_line.encode('latin-1', 'replace').decode('latin-1')) # Basic non-utf8 handling
            except Exception as e:
                logger.warning(f"Skipping problematic line: {raw_line[:50]}... ({e})")
                continue
            pdf.ln(2) # Add a bit of space between paragraphs

        pdf.output(filepath)
        logger.info(f"✅ PDF generated successfully: {filepath}")

        return {"status": "success", "result": filepath, "file_size": os.path.getsize(filepath)}

    except Exception as e:
        logger.error(f"pdf_generator exception: {e}", exc_info=True)
        return {"status": "error", "message": f"Error: PDF generation failed."}

# --- NEW TOOL ---
@tool
def presentation_generator(topic: str, content: str, filename: str) -> str:
    """
    Generate a simple .pptx presentation file from content.
    
    Args:
        inputs (dict): A dictionary containing:
            - "topic" (str): The main topic for the title slide.
            - "content" (str): The report or text content, split by newlines for slides.
            - "filename" (str): The desired output filename (e.g., "report.pptx").
            
    Returns:
        JSON string with status and the filepath, or error message.
    """
    try:
        # topic = inputs.get("topic", "Untitled Presentation") # No longer need this
        # content = inputs.get("content", "No content provided.") # No longer need this
        # filename = inputs.get("filename", "presentation.pptx") # No longer need this

        if not filename.endswith(".pptx"):
            filename += ".pptx"

        script_dir = os.path.dirname(os.path.abspath(__file__))
        output_dir = os.path.abspath(os.path.join(script_dir, "..", "generated_documents"))

        os.makedirs(output_dir, exist_ok=True)
        filepath = os.path.join(output_dir, filename)

        prs = Presentation()
        
        # Title slide
        title_slide_layout = prs.slide_layouts[0]
        slide = prs.slides.add_slide(title_slide_layout)
        title_placeholder = slide.shapes.title
        subtitle_placeholder = slide.placeholders[1]
        title_placeholder.text = topic
        subtitle_placeholder.text = "Generated by Educational Agent"

        # Content slides
        content_slide_layout = prs.slide_layouts[1] # Title and Content
        
        # Split content into rough slides (e.g., by paragraph)
        sections = [s.strip() for s in content.split("\n\n") if s.strip()]
        
        if not sections:
            sections = [content] # Use the whole content if no double-newlines

        for section in sections[:9]: # Limit to 9 content slides
            slide = prs.slides.add_slide(content_slide_layout)
            title_placeholder = slide.shapes.title
            body_placeholder = slide.placeholders[1]

            # Try to find a 'title' for the slide (first line)
            lines = section.split('\n')
            title_placeholder.text = lines[0][:80] # Use first line as title, truncated
            
            # Use rest of lines as bullet points
            if len(lines) > 1:
                tf = body_placeholder.text_frame
                tf.clear() # Clear existing placeholder text
                p = tf.paragraphs[0]
                p.text = lines[1]
                p.font.size = Pt(18)
                
                for line in lines[2:]:
                    p = tf.add_paragraph()
                    p.text = line
                    p.font.size = Pt(18)
                    p.level = 0 # Top-level bullet

        prs.save(filepath)
        logger.info(f"✅ Presentation generated successfully: {filepath}")

        return json.dumps({"status": "success", "result": filepath, "file_size": os.path.getsize(filepath)})

    except Exception as e:
        logger.error(f"presentation_generator exception: {e}", exc_info=True)
        return json.dumps({"status": "error", "message": f"Error: Presentation generation failed."})


@tool
def email_tool(to_email: str, subject: str, body: str, attachment_paths: list[str] | None = None) -> str:
    """
    Send email with optional attachments.
    
    Args:
        to_email: Recipient email address
        subject: Email subject line
        body: Email body text
        attachment_paths: Optional list of file paths to attach
        
    Returns:
        JSON string with status and result, or error message
    """
    try:
        if not to_email or not subject or not body:
            # ... (error handling) ...
            pass
        
        smtp_server = os.getenv("EMAIL_SMTP_SERVER", "smtp.gmail.com")
        smtp_port = int(os.getenv("EMAIL_SMTP_PORT", "587"))
        sender_email = os.getenv("EMAIL_USERNAME", "periodictablesme@gmail.com")
        sender_password = os.getenv("EMAIL_PASSWORD", "rjxl wmay mmuk yxav")
        
        if not sender_email or not sender_password:
            # ... (error handling) ...
            pass
        
        msg = MIMEMultipart()
        msg["From"] = sender_email
        msg["To"] = to_email
        msg["Subject"] = subject
        msg.attach(MIMEText(body, "plain"))
        
        attachment_info = ""
        
        if attachment_paths:
            for path_str in attachment_paths:
                attachment_path = Path(path_str)
                if not attachment_path.exists():
                    logger.warning(f"email_tool: Attachment file not found, skipping: {attachment_path}")
                    continue
                
                try:
                    with open(attachment_path, "rb") as attachment:
                        part = MIMEBase("application", "octet-stream")
                        part.set_payload(attachment.read())
                    encoders.encode_base64(part)
                    part.add_header("Content-Disposition", f"attachment; filename= {attachment_path.name}")
                    msg.attach(part)
                    attachment_info += f" {attachment_path.name}"
                except Exception as e:
                    logger.error(f"email_tool attachment error: {e}", exc_info=True)
                    # Continue trying to send email without this attachment
            
            if attachment_info:
                attachment_info = f" with attachments:{attachment_info}"

        # ... (SMTP send logic) ...
        try:
            server = smtplib.SMTP(smtp_server, smtp_port)
            server.starttls()
            server.login(sender_email, sender_password)
            server.sendmail(sender_email, to_email, msg.as_string())
            server.quit()
            
            success_msg = f"Email sent to {to_email}{attachment_info}"
            logger.info(f"email_tool: {success_msg}")
            
            return json.dumps({"status": "success", "result": success_msg})
            
        except smtplib.SMTPException as e:
            logger.error(f"email_tool SMTP error: {e}", exc_info=True)
            return json.dumps({"status": "error", "message": "Error: SMTP error occurred."})
        
    except Exception as e:
        logger.error(f"email_tool exception: {e}", exc_info=True)
        return json.dumps({"status": "error", "message": "Error: Email sending failed."})

# ======================== EXPORTS ========================

__all__ = [
    "knowledge_retrieval",
    "quiz_generator",
    "pdf_generator",
    "email_tool",
    "report_generator",       # NEW
    "presentation_generator"  # NEW
]