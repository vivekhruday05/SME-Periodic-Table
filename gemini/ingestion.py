"""
Preprocessing & Chunking pipeline

Saves chunked documents as JSONL with metadata and parent-child relationships.

Features:
- Extract text from PDF, DOCX, PPTX, TXT, MD, CSV (uses pdfplumber, python-docx, python-pptx if available)
- [BONUS] OCR fallback for scanned PDFs (requires pytesseract)
- Cleans text: lowercasing (optional), removes non-informative content (headers/footers heuristics), normalizes whitespace
- Deduplication by SHA256 of cleaned text
- Multiple-granularity chunking: 2048, 512, 128 tokens (configurable)
- Chunking strategies: fixed-token (sliding) and paragraph-aware packing (default)
- Context-aware overlap (configurable overlap fraction)
- [BONUS] Batch ingestion pipeline with multiprocessing for faster extraction
- Comprehensive logging and error handling
- Outputs:
  - data/processed/chunks.jsonl  (one line per chunk, includes parent-child metadata)
  - data/processed/docs_manifest.jsonl (one line per source document metadata)
  - logs/ingest.log and logs/errors.log

Chunking Strategy Justification:
The default strategy is 'paragraph'. This is a content-aware approach that
treats paragraphs (separated by double newlines) as the atomic unit. It
tries to pack as many whole paragraphs as possible into each chunk without
exceeding the token limit. This is superior to 'sliding' (fixed-size)
because it avoids splitting sentences or paragraphs mid-thought, preserving
the semantic integrity of the text within each chunk. This leads to
better context for retrieval and downstream LLM tasks. The 'sliding'
strategy is provided as a fallback for text without clear paragraph breaks.

Usage:
    python utils/preprocessing_chunking_pipeline.py --data-root ./data --out-dir ./data/processed --batch-size 8 --enable-ocr

Dependencies (optional but recommended):
    pip install pdfplumber python-docx python-pptx tqdm
    pip install pytesseract pillow (for OCR)
"""

import os
import re
import sys
import json
import uuid
import hashlib
import argparse
import logging
import multiprocessing
import csv
from pathlib import Path
from typing import List, Dict, Tuple
from datetime import datetime

# Optional imports
try:
    import pdfplumber
except Exception:
    pdfplumber = None

try:
    import pytesseract
    from PIL import Image
except Exception:
    pytesseract = None
    Image = None

try:
    import docx
except Exception:
    docx = None

try:
    from pptx import Presentation
except Exception:
    Presentation = None

try:
    from tqdm import tqdm
except Exception:
    tqdm = lambda x, **kw: x  # No-op if tqdm isn't installed

# ------------------- Logging setup -------------------

def setup_logging(log_dir: Path):
    """Sets up primary and error loggers."""
    log_dir.mkdir(parents=True, exist_ok=True)
    ingest_log = log_dir / 'ingest.log'
    error_log = log_dir / 'errors.log'

    # Configure the 'ingest' logger
    logger = logging.getLogger('ingest')
    logger.setLevel(logging.INFO)
    # Prevent logs from propagating to the root logger
    logger.propagate = False

    # (Re)configure handlers if they already exist from a previous setup
    if not logger.handlers:
        fh = logging.FileHandler(ingest_log, encoding='utf-8')
        fh.setLevel(logging.INFO)
        ch = logging.StreamHandler(sys.stdout)
        ch.setLevel(logging.INFO)

        formatter = logging.Formatter('%(asctime)s - %(processName)s - %(levelname)s - %(message)s')
        fh.setFormatter(formatter)
        ch.setFormatter(formatter)
    
        logger.addHandler(fh)
        logger.addHandler(ch)

    # Configure the 'errors' logger
    err_logger = logging.getLogger('errors')
    err_logger.setLevel(logging.WARNING)
    err_logger.propagate = False

    if not err_logger.handlers:
        efh = logging.FileHandler(error_log, encoding='utf-8')
        efh.setLevel(logging.WARNING)
        e_formatter = logging.Formatter('%(asctime)s - %(processName)s - %(levelname)s - %(message)s')
        efh.setFormatter(e_formatter)
        err_logger.addHandler(efh)

    return logger, err_logger

# ------------------- Multiprocessing Worker Setup -------------------

# Global dict to hold loggers for worker processes
worker_loggers = {}

def init_worker_logging(log_dir: Path):
    """Initializer function for multiprocessing pool to set up logging."""
    logger, err_logger = setup_logging(log_dir)
    worker_loggers['ingest'] = logger
    worker_loggers['errors'] = err_logger

# ------------------- Utilities -------------------

def sha256_text(text: str) -> str:
    """Computes SHA256 hash of a string."""
    return hashlib.sha256(text.encode('utf-8')).hexdigest()

# ------------------- Text extractors -------------------

def extract_text_pdf(path: Path, enable_ocr: bool = False, logger=None) -> str:
    """Extracts text from PDF, with optional OCR fallback."""
    if pdfplumber is None:
        if logger:
            logger.warning('pdfplumber not installed. Skipping PDF extraction.')
        return ""
        
    text_parts = []
    try:
        with pdfplumber.open(path) as pdf:
            for i, page in enumerate(pdf.pages):
                page_text = page.extract_text() or ''
                
                # Check if text is minimal or absent, and if OCR is enabled
                if (not page_text.strip() or len(page_text.strip()) < 20) and enable_ocr:
                    if pytesseract is None or Image is None:
                        if logger and i == 0: # Log only once per file
                             logger.warning(f"OCR enabled but pytesseract/Pillow not installed. OCR will be skipped for {path.name}.")
                        enable_ocr = False # Disable for this run
                    else:
                        try:
                            # Use a reasonable resolution for OCR
                            image = page.to_image(resolution=200).original
                            ocr_text = pytesseract.image_to_string(image)
                            if ocr_text and ocr_text.strip():
                                if logger:
                                    logger.info(f"OCR used for page {i+1} of {path.name}")
                                page_text = ocr_text
                        except Exception as e:
                            if logger:
                                logger.warning(f"OCR failed for page {i+1} of {path.name}: {e}")
                                
                text_parts.append(page_text)
    except Exception as e:
        if logger:
            logger.error(f"Failed to open PDF {path.name}: {e}")
        return "" # Return empty string on failure to open
        
    return '\n\n'.join(text_parts) # Join pages with double newline for paragraph separation


def extract_text_docx(path: Path, logger=None) -> str:
    """Extracts text from DOCX."""
    if docx is None:
        if logger:
            logger.warning('python-docx not installed. Skipping DOCX extraction.')
        return ""
    try:
        doc = docx.Document(path)
        return '\n\n'.join(p.text for p in doc.paragraphs if p.text.strip())
    except Exception as e:
        if logger:
            logger.error(f"Failed to read DOCX {path.name}: {e}")
        return ""


def extract_text_pptx(path: Path, logger=None) -> str:
    """Extracts text from PPTX."""
    if Presentation is None:
        if logger:
            logger.warning('python-pptx not installed. Skipping PPTX extraction.')
        return ""
    try:
        prs = Presentation(str(path))
        lines = []
        for slide in prs.slides:
            slide_text = []
            for shape in slide.shapes:
                if hasattr(shape, 'text'):
                    slide_text.append(shape.text)
            lines.append('\n'.join(slide_text))
        return '\n\n'.join(lines) # Join slides with double newline
    except Exception as e:
        if logger:
            logger.error(f"Failed to read PPTX {path.name}: {e}")
        return ""


def extract_text_txt_md(path: Path, logger=None) -> str:
    """Extracts text from TXT or MD."""
    try:
        with open(path, 'r', encoding='utf-8') as f:
            return f.read()
    except UnicodeDecodeError:
        try:
            with open(path, 'r', encoding='latin-1') as f:
                return f.read()
        except Exception as e:
            if logger:
                logger.error(f"Failed to read text file {path.name} (even as latin-1): {e}")
            return ""
    except Exception as e:
        if logger:
             logger.error(f"Failed to read text file {path.name}: {e}")
        return ""

def extract_text_csv(path: Path, logger=None) -> str:
    """Extracts text from CSV by reading rows and joining them line by line.
    Uses csv.Sniffer to detect dialect when possible. Falls back gracefully on errors.
    """
    def _read_csv(encoding: str) -> str:
        try:
            with open(path, 'r', encoding=encoding, newline='') as f:
                sample = f.read(4096)
                f.seek(0)
                try:
                    dialect = csv.Sniffer().sniff(sample)
                except Exception:
                    dialect = csv.excel
                reader = csv.reader(f, dialect)
                lines = []
                for row in reader:
                    # Join fields with a tab for readability
                    safe_fields = [(col or '').strip() for col in row]
                    lines.append('\t'.join(safe_fields))
                return '\n'.join(lines)
        except Exception as e:
            if logger:
                logger.error(f"Failed to read CSV {path.name} with encoding {encoding}: {e}")
            raise

    try:
        return _read_csv('utf-8')
    except Exception:
        try:
            return _read_csv('latin-1')
        except Exception:
            return ""

def extract_text_generic(path: Path, enable_ocr=False, logger=None) -> str:
    """Generic dispatcher for text extraction based on file extension."""
    suffix = path.suffix.lower()
    if suffix == '.pdf':
        return extract_text_pdf(path, enable_ocr=enable_ocr, logger=logger)
    elif suffix == '.docx':
        return extract_text_docx(path, logger=logger)
    elif suffix == '.pptx':
        return extract_text_pptx(path, logger=logger)
    elif suffix == '.csv':
        return extract_text_csv(path, logger=logger)
    elif suffix in ('.txt', '.md', '.py', '.js', '.css', '.html', '.json'):
        return extract_text_txt_md(path, logger=logger)
    else:
        if logger:
            logger.info(f"Unsupported file type {suffix} for {path.name}. Attempting plain text read.")
        # Fallback for unknown but text-like files
        return extract_text_txt_md(path, logger=logger)

# ------------------- Cleaning & Preprocessing -------------------

def normalize_whitespace(text: str) -> str:
    """Cleans and normalizes whitespace in text."""
    text = text.replace('\r', '\n')
    text = re.sub(r'[ \t]+', ' ', text) # collapse horizontal whitespace
    text = re.sub(r'\n{3,}', '\n\n', text) # collapse >2 newlines to 2 (paragraph break)
    text = re.sub(r' \n', '\n', text) # remove spaces before newlines
    text = re.sub(r'\n ', '\n', text) # remove spaces after newlines
    # Re-join lines to strip trailing spaces, then rejoin
    text = '\n'.join([ln.strip() for ln in text.splitlines()])
    return text.strip()


def remove_non_informative(text: str, min_sentence_len: int = 3) -> str:
    """Heuristically removes headers, footers, and other non-informative lines."""
    lines = []
    
    # Regex to match lines that are *only* a date, time, or timestamp
    # This is intentionally strict to avoid matching full sentences.
    date_regex = r'(\d{1,2}[/-]\d{1,2}[/-]\d{2,4})' # e.g., 11/02/2025 or 11-02-25
    time_regex = r'(\d{1,2}:\d{2}(:\d{2})?(\s?[AP]M)?)' # e.g., 10:51:57 PM or 10:51
    text_date_regex = r'((Jan|Feb|Mar|Apr|May|Jun|Jul|Aug|Sep|Oct|Nov|Dec)\s\d{1,2},\s\d{4})' # e.g., Nov 2, 2025
    
    # Combine them into a single pattern for re.fullmatch
    # Allows for optional surrounding whitespace
    standalone_datetime_pattern = re.compile(
        rf'\s*({date_regex}|{time_regex}|{text_date_regex})\s*', 
        re.IGNORECASE
    )

    for ln in text.splitlines():
        stripped = ln.strip()
        if len(stripped) == 0:
            continue
        
        # skip lines that look like page numbers (e.g., "123", "Page 12")
        if re.fullmatch(r'(Page\s+)?\d{1,4}(\s+of\s+\d{1,4})?', stripped, re.IGNORECASE):
            continue
            
        # --- NEW RULE ---
        # skip lines that are *just* a date or time
        if re.fullmatch(standalone_datetime_pattern, stripped):
            continue
        # --- END NEW RULE ---
            
        # skip lines that are mostly punctuation or symbols
        non_alnum = re.sub(r'[\w\s]', '', stripped)
        if len(non_alnum) > (len(stripped) * 0.6) and len(stripped) > 5:
             continue # e.g., "**********" or "---...---"

        # skip very short lines that aren't part of a list
        if len(stripped.split()) < min_sentence_len:
            # allow short lines if they look like list items or headings
            if not re.match(r'^(\*|\-|\d+\.|\(?[a-zA-Z\d]\))', stripped):
                continue
                
        lines.append(stripped)
    return '\n'.join(lines)

def clean_text(text: str, lowercase: bool = True) -> str:
    """Applies all cleaning steps to the extracted text."""
    
    # 1. Apply line-level filtering first
    text = remove_non_informative(text)
    
    # 2. Normalize whitespace (this is a good time to do it)
    text = normalize_whitespace(text)

    # --- ADD/UPDATE CLEANING RULES ---
    
    # 3. Collapse "stuttering" character repetitions from bad OCR/copy-paste
    # This is the fix for 'nnnnnooooommmm' -> 'nome'
    # It finds any character repeated 3 or more times and replaces with a single one.
    text = re.sub(r'(.)\1{2,}', r'\1', text)
    
    # 4. Remove "word art" headings (e.g., "n o m e n c l a t u r e")
    # This finds sequences of 5 or more single letters followed by a space
    text = re.sub(r'(?:[a-zA-Z]\s){5,}', ' ', text)
    
    # 5. Remove long standalone noise (e.g., "55555" or "...." or "____")
    # This removes any "word" of 4+ digits, dots, underscores, or dashes
    text = re.sub(r'\b([\d._\-]){4,}\b', ' ', text)
    
    # 6. Remove any remaining long strings of just punctuation
    text = re.sub(r'[.…*#=_-]{4,}', ' ', text)

    # --- END OF NEW RULES ---

    if lowercase:
        text = text.lower()
        
    # Replace common unicode oddities
    text = text.replace('\u2022', '*') # bullet
    text = text.replace('\u201c', '"').replace('\u201d', '"') # smart quotes
    text = text.replace('\u2018', "'").replace('\u2019', "'") # smart quotes
    text = text.replace('\u2013', '-').replace('\u2014', '--') # dashes
    
    # 7. Apply normalization one last time to clean up spaces left by regex
    text = normalize_whitespace(text)
    
    return text.strip()
# ------------------- Tokenization (estimator) -------------------

def whitespace_tokenize(text: str) -> List[str]:
    """
    Simple, fast whitespace-based tokenization.
    Used for estimating chunk size. Not a true tokenizer.
    """
    tokens = []
    for chunk in re.split(r'\s+', text):
        if chunk:
            tokens.append(chunk)
    return tokens

# ------------------- Chunking strategies -------------------

def pack_paragraphs_into_chunks(paragraphs: List[str], chunk_size: int, overlap: int, tokenizer_fn) -> List[Tuple[str,int,int]]:
    """
    Packs paragraphs into chunks, respecting chunk_size.
    Overlap is implemented by re-including previous paragraphs.
    Returns list of tuples: (chunk_text, start_para_index, end_para_index)
    """
    para_tokens = [(p, len(tokenizer_fn(p))) for p in paragraphs]
    chunks = []

    i = 0
    n = len(paragraphs)
    while i < n:
        cur_tokens = 0
        j = i
        parts = []
        
        # Accumulate paragraphs until chunk_size is reached
        while j < n and (cur_tokens + para_tokens[j][1]) <= chunk_size:
            parts.append(para_tokens[j][0])
            cur_tokens += para_tokens[j][1]
            j += 1
            
        # If a single paragraph is larger than chunk_size, split it
        if i == j:
            # Force-split the long paragraph by tokens
            long_para_text = para_tokens[i][0]
            long_para_tokens = tokenizer_fn(long_para_text)
            
            start = 0
            while start < len(long_para_tokens):
                end = min(start + chunk_size, len(long_para_tokens))
                chunk_text = ' '.join(long_para_tokens[start:end])
                chunks.append((chunk_text, i, i))
                
                new_start = end - overlap
                # Ensure we make progress and avoid infinite loops
                if new_start <= start:
                    start = end
                else:
                    start = new_start
            i += 1
            continue

        # We have a valid chunk
        chunk_text = '\n\n'.join(parts)
        chunks.append((chunk_text, i, j-1)) # j-1 is the last included index
        
        # Calculate overlap
        if overlap <= 0 or j >= n:
            i = j # No overlap or end of document
        else:
            # Step back from j-1 to accumulate 'overlap' tokens
            overlap_token_count = 0
            k = j - 1
            while k > i and overlap_token_count < overlap:
                overlap_token_count += para_tokens[k][1]
                k -= 1
            
            # New start index 'i' is k+1 (or j-1 if overlap is small)
            i = k + 1
            # Ensure we always make progress
            if i >= j:
                i = j 

    return chunks


def sliding_window_token_chunks(text: str, chunk_size: int, overlap: int, tokenizer_fn) -> List[Tuple[str,int,int]]:
    """
    Simple fixed-size sliding window chunking.
    Returns list of tuples: (chunk_text, start_token_index, end_token_index)
    """
    tokens = tokenizer_fn(text)
    n = len(tokens)
    if n == 0:
        return []
        
    chunks = []
    start = 0
    while start < n:
        end = min(start + chunk_size, n)
        chunk_text = ' '.join(tokens[start:end])
        # Using token indices for start/end position
        chunks.append((chunk_text, start, end))
        
        if end == n:
            break # Reached the end
            
        new_start = end - overlap
        # Ensure we make progress
        if new_start <= start:
            start = start + 1 
        else:
            start = new_start
            
    return chunks

# ------------------- Deduplication -------------------

def deduplicate_documents(docs: List[Dict]) -> Tuple[List[Dict], List[Dict]]:
    """Deduplicates a list of document metadata dicts based on 'sha256' key."""
    seen = set()
    unique = []
    dupes = []
    for doc in docs:
        h = doc['sha256']
        if not h: # Skip docs that failed extraction
            dupes.append(doc)
            continue
            
        if h in seen:
            dupes.append(doc)
        else:
            seen.add(h)
            unique.append(doc)
    return unique, dupes

# ------------------- Main pipeline worker -------------------

def process_single_file_worker(path: Path, doc_id: str, lowercase: bool, enable_ocr: bool, base_meta: Dict = None) -> Dict:
    """
    Worker function to process a single file.
    Fetches loggers from the global 'worker_loggers' dict.
    Accepts an optional base_meta dictionary.
    """
    # Get loggers initialized by the pool
    logger = worker_loggers.get('ingest')
    err_logger = worker_loggers.get('errors')

    meta = base_meta or {}
    meta.update({
        'id': doc_id,
        'source_path': str(path),
        'filename': path.name,
        'mtime': datetime.fromtimestamp(path.stat().st_mtime).isoformat(),
        'size_bytes': path.stat().st_size,
    })
    try:
        # ** OCR FIX IS HERE: passing enable_ocr and loggers **
        text = extract_text_generic(path, enable_ocr=enable_ocr, logger=logger)
        
        if not text or len(text.strip()) == 0:
            raise RuntimeError('No text extracted')
            
        cleaned = clean_text(text, lowercase=lowercase)
        
        if not cleaned or len(cleaned.strip()) == 0:
            raise RuntimeError('Text was empty after cleaning')

        meta['cleaned_text_snippet'] = cleaned[:1000] # Larger snippet for review
        meta['sha256'] = sha256_text(cleaned)
        meta['word_count_est'] = len(whitespace_tokenize(cleaned))
        meta['extracted'] = True
        
        if logger:
            logger.info(f"Extracted and cleaned: {path.name} ({meta['word_count_est']} tokens)")
            
    except Exception as e:
        if err_logger:
            err_logger.warning(f"Failed to extract {path}: {e}")
            
        meta['extracted'] = False
        meta['error'] = str(e)
        meta['sha256'] = None
        meta['word_count_est'] = 0
        meta['cleaned_text_snippet'] = ''
        
    return meta

# ------------------- Main pipeline controller -------------------

def ingest_directory(data_root: Path, out_dir: Path,
                     chunk_sizes: Tuple[int] = (2048, 512, 128),
                     overlap_frac: float = 0.1,
                     split_strategy: str = 'paragraph',
                     lowercase: bool = True,
                     batch_size: int = 4,
                     enable_ocr: bool = False,
                     tokenizer_fn=whitespace_tokenize):
    """
    Main function to run the ingestion and chunking pipeline.
    """
    log_dir = out_dir.parent / 'logs'
    logger, err_logger = setup_logging(log_dir)

    # 0. Load external metadata if it exists
    external_metadata_map = {}
    metadata_file = data_root / 'metadata.jsonl'
    if metadata_file.is_file():
        logger.info(f"Loading external metadata from {metadata_file}")
        with open(metadata_file, 'r', encoding='utf-8') as f:
            for line in f:
                try:
                    entry = json.loads(line)
                    if 'filepath' in entry:
                        # The key is the filepath relative to the data_root
                        key_path = entry['filepath']
                        external_metadata_map[key_path] = entry
                except (json.JSONDecodeError, KeyError) as e:
                    err_logger.warning(f"Skipping malformed line in {metadata_file}: {e}")
        logger.info(f"Loaded {len(external_metadata_map)} external metadata entries.")


    # 1. Find files
    files = [p for p in data_root.rglob('*') 
             if p.is_file() and 
             p.name.lower() not in ('metadata.jsonl', '.ds_store') and
             not p.name.startswith('~')]
             
    logger.info(f'Found {len(files)} files under {data_root}')
    if not files:
        logger.warning("No files found to process.")
        return None, None

    # 2. **BATCH PROCESSING (FIXED) **
    # Use multiprocessing.Pool to process files in parallel.
    # 'batch_size' is used as the number of worker processes.
    num_workers = max(1, min(batch_size, os.cpu_count() or 1))
    logger.info(f"Starting batch processing with {num_workers} workers...")
    
    # Create task list for starmap
    tasks = []
    for p in files:
        relative_path_str = str(p.relative_to(data_root))
        base_meta = external_metadata_map.get(relative_path_str, {})
        tasks.append((p, str(uuid.uuid4()), lowercase, enable_ocr, base_meta))

    docs_meta = []

    try:
        # 'init_worker_logging' is called once per worker process
        with multiprocessing.Pool(processes=num_workers, 
                                  initializer=init_worker_logging, 
                                  initargs=(log_dir,)) as pool:
            
            # Use tqdm to show progress for the parallel processing
            with tqdm(total=len(tasks), desc='Processing files') as pbar:
                # starmap applies the function to each tuple in 'tasks'
                for meta in pool.starmap(process_single_file_worker, tasks):
                    docs_meta.append(meta)
                    pbar.update()
                    
    except Exception as e:
        logger.error(f"Critical error during parallel processing: {e}")
        return None, None

    # 3. Deduplicate
    unique_docs, dupes = deduplicate_documents(docs_meta)
    logger.info(f'Processed {len(docs_meta)} files. Found {len(unique_docs)} unique, {len(dupes)} duplicates/errors.')
    if dupes:
        for d in dupes:
            if d.get('error'):
                err_logger.warning(f"File failed processing: {d['filename']} | Error: {d['error']}")
            else:
                logger.info(f"Duplicate file found (will be skipped): {d['filename']} (SHA: {d['sha256']})")


    # 4. Write manifest
    out_dir.mkdir(parents=True, exist_ok=True)
    manifest_path = out_dir / 'docs_manifest.jsonl'
    with open(manifest_path, 'w', encoding='utf-8') as mf:
        for d in unique_docs:
            mf.write(json.dumps(d, ensure_ascii=False) + '\n')
    logger.info(f"Wrote manifest for {len(unique_docs)} unique docs to {manifest_path}")

    # 5. Chunking (now hierarchical)
    chunks_out = out_dir / 'chunks.jsonl'
    total_chunks = 0
    
    # Ensure chunk sizes are sorted from largest to smallest
    chunk_sizes = sorted(chunk_sizes, reverse=True)

    with open(chunks_out, 'w', encoding='utf-8') as cf:
        for d in tqdm(unique_docs, desc='Chunking docs'):
            if not d.get('extracted', False):
                continue
            
            doc_parent_id = d['id']
            
            try:
                raw_text = extract_text_generic(Path(d['source_path']), enable_ocr=enable_ocr, logger=logger)
                cleaned_full = clean_text(raw_text, lowercase=lowercase)
                if not cleaned_full:
                    continue
            except Exception as e:
                err_logger.warning(f"Failed to re-extract for chunking {d['source_path']}: {e}")
                continue

            # Start the recursive chunking process
            # The top-level chunks have the document as their parent
            
            # Base case: chunk the full document into the largest chunk size
            largest_chunk_size = chunk_sizes[0]
            overlap_tokens = max(1, int(largest_chunk_size * overlap_frac))
            
            paragraphs = [p for p in re.split(r'\n{2,}', cleaned_full) if p.strip()]
            
            if split_strategy == 'paragraph' and len(paragraphs) > 0:
                chunk_tuples = pack_paragraphs_into_chunks(paragraphs, largest_chunk_size, overlap_tokens, tokenizer_fn)
            else:
                if split_strategy == 'paragraph':
                     logger.info(f"No paragraphs found in {d['filename']}, falling back to 'sliding' strategy for top level.")
                chunk_tuples = sliding_window_token_chunks(cleaned_full, largest_chunk_size, overlap_tokens, tokenizer_fn)

            
            # This list will hold all chunks of all sizes for this document
            all_doc_chunks = []

            # Process the largest chunks and then recursively chunk them
            for idx, (chunk_text, s, e) in enumerate(chunk_tuples):
                
                # This is a top-level (largest) chunk
                parent_chunk_id = str(uuid.uuid4())
                chunk_entry = {
                    'chunk_id': parent_chunk_id,
                    'parent_id': doc_parent_id, # Document is the parent
                    'parent_chunk_id': None, # No parent chunk
                    'source_path': d['source_path'],
                    'filename': d['filename'],
                    'chunk_size_tokens': largest_chunk_size,
                    'overlap_tokens': overlap_tokens,
                    'chunk_index_for_size': idx,
                    'chunk_text': chunk_text,
                    'start_position': s,
                    'end_position': e,
                    'chunk_word_count_est': len(tokenizer_fn(chunk_text)),
                    'created_at': datetime.utcnow().isoformat() + 'Z'
                }
                # Carry over external metadata
                for key, value in d.items():
                    if key not in chunk_entry and key not in ['id', 'cleaned_text_snippet', 'sha256', 'word_count_est', 'extracted', 'error', 'mtime', 'size_bytes']:
                        chunk_entry[key] = value
                
                all_doc_chunks.append(chunk_entry)

                # Now, recursively chunk this large chunk into smaller ones
                child_chunks = []
                parent_text_for_children = chunk_text
                
                for i in range(1, len(chunk_sizes)):
                    child_chunk_size = chunk_sizes[i]
                    parent_chunk_size = chunk_sizes[i-1]
                    
                    # The text to be chunked is from the parent chunk
                    text_to_chunk = parent_text_for_children
                    
                    child_overlap = max(1, int(child_chunk_size * overlap_frac))
                    
                    # We use sliding window for sub-chunking as paragraph context is already coarse
                    child_tuples = sliding_window_token_chunks(text_to_chunk, child_chunk_size, child_overlap, tokenizer_fn)
                    
                    # The parent for these new chunks is the one we just created
                    # If we are creating 128-token chunks from a 512, the 512 is the parent.
                    # We need to find the right parent from the previous level of chunks.
                    # This simplified logic assumes we are chunking the *first* big chunk.
                    # A better logic would be needed for multiple big chunks.
                    # For this implementation, we'll just pass the current parent_chunk_id down.
                    
                    new_children = []
                    for child_idx, (child_text, cs, ce) in enumerate(child_tuples):
                        child_id = str(uuid.uuid4())
                        child_entry = {
                            'chunk_id': child_id,
                            'parent_id': doc_parent_id,
                            'parent_chunk_id': parent_chunk_id, # Link to immediate parent chunk
                            'source_path': d['source_path'],
                            'filename': d['filename'],
                            'chunk_size_tokens': child_chunk_size,
                            'overlap_tokens': child_overlap,
                            'chunk_index_for_size': child_idx,
                            'chunk_text': child_text,
                            'start_position': cs,
                            'end_position': ce,
                            'chunk_word_count_est': len(tokenizer_fn(child_text)),
                            'created_at': datetime.utcnow().isoformat() + 'Z'
                        }
                        # Carry over metadata
                        for key, value in d.items():
                            if key not in child_entry and key not in ['id', 'cleaned_text_snippet', 'sha256', 'word_count_est', 'extracted', 'error', 'mtime', 'size_bytes']:
                                child_entry[key] = value
                        
                        new_children.append(child_entry)
                    
                    child_chunks.extend(new_children)
                    
                    # For the next level of chunking, the parent_chunk_id needs to be updated
                    # This part is tricky. Let's assume for now we just keep adding them.
                    # The parent_chunk_id should be the ID of the chunk from which it was created.
                    # The current logic is flawed as it doesn't create a clear hierarchy beyond one level.
                    
                all_doc_chunks.extend(child_chunks)


            # A more robust recursive implementation
            def create_hierarchical_chunks(text_to_chunk, parent_id, parent_is_doc, level):
                if level >= len(chunk_sizes):
                    return []

                current_chunk_size = chunk_sizes[level]
                overlap = max(1, int(current_chunk_size * overlap_frac))
                
                # Use paragraph strategy only for the top level (doc -> large chunks)
                if parent_is_doc:
                    paragraphs = [p for p in re.split(r'\n{2,}', text_to_chunk) if p.strip()]
                    if split_strategy == 'paragraph' and len(paragraphs) > 0:
                        tuples = pack_paragraphs_into_chunks(paragraphs, current_chunk_size, overlap, tokenizer_fn)
                    else:
                        tuples = sliding_window_token_chunks(text_to_chunk, current_chunk_size, overlap, tokenizer_fn)
                else: # Sub-chunking is always sliding window
                    tuples = sliding_window_token_chunks(text_to_chunk, current_chunk_size, overlap, tokenizer_fn)

                
                created_chunks = []
                for idx, (chunk_text, s, e) in enumerate(tuples):
                    chunk_id = str(uuid.uuid4())
                    
                    entry = {
                        'chunk_id': chunk_id,
                        'parent_id': doc_parent_id, # Always the doc ID
                        'parent_chunk_id': None if parent_is_doc else parent_id,
                        'source_path': d['source_path'],
                        'filename': d['filename'],
                        'chunk_size_tokens': current_chunk_size,
                        'overlap_tokens': overlap,
                        'chunk_index_for_size': idx,
                        'chunk_text': chunk_text,
                        'start_position': s,
                        'end_position': e,
                        'chunk_word_count_est': len(tokenizer_fn(chunk_text)),
                        'created_at': datetime.utcnow().isoformat() + 'Z'
                    }
                    # Carry over metadata
                    for key, value in d.items():
                        if key not in entry and key not in ['id', 'cleaned_text_snippet', 'sha256', 'word_count_est', 'extracted', 'error', 'mtime', 'size_bytes']:
                            entry[key] = value
                    
                    created_chunks.append(entry)
                    
                    # Recurse to create children chunks
                    children = create_hierarchical_chunks(chunk_text, chunk_id, False, level + 1)
                    created_chunks.extend(children)
                    
                return created_chunks

            # Start the process for the whole document
            final_chunks = create_hierarchical_chunks(cleaned_full, doc_parent_id, True, 0)

            for chunk in final_chunks:
                cf.write(json.dumps(chunk, ensure_ascii=False) + '\n')
                total_chunks += 1
                    
    logger.info(f'Wrote {total_chunks} chunks to {chunks_out}')
    return manifest_path, chunks_out

# ------------------- CLI -------------------

def main(argv=None):
    parser = argparse.ArgumentParser(description='Preprocessing & Chunking Pipeline with OCR fallback')
    parser.add_argument('--data-root', type=str, default='./data', help='Directory containing raw documents.')
    parser.add_argument('--out-dir', type=str, default='./processed', help='Directory to save processed JSONL files.')
    parser.add_argument('--chunk-sizes', type=str, default='2048,512,128', help='Comma-separated list of token chunk sizes.')
    parser.add_argument('--overlap-frac', type=float, default=0.1, help='Fraction of chunk_size to use as overlap (in tokens).')
    parser.add_argument('--split-strategy', type=str, default='paragraph', choices=['paragraph', 'sliding'], help='Chunking strategy.')
    parser.add_argument('--no-lowercase', action='store_true', help='Disable lowercasing of text.')
    parser.add_argument('--enable-ocr', action='store_true', help='[BONUS] Enable OCR fallback for image-based PDFs (requires pytesseract).')
    parser.add_argument('--batch-size', type=int, default=4, help='[BONUS] Number of parallel processes for file extraction.')
    
    args = parser.parse_args(argv)

    data_root = Path(args.data_root)
    out_dir = Path(args.out_dir)
    chunk_sizes = [int(x.strip()) for x in args.chunk_sizes.split(',') if x.strip()]
    overlap_frac = args.overlap_frac
    lowercase = not args.no_lowercase

    if args.enable_ocr and (pytesseract is None or Image is None):
        print("WARNING: --enable-ocr was specified, but 'pytesseract' or 'Pillow' is not installed. OCR will be skipped.", file=sys.stderr)
        print("Please run: pip install pytesseract pillow", file=sys.stderr)

    # ** FIX: Pass batch_size and enable_ocr to the main function **
    ingest_directory(data_root, out_dir, 
                     chunk_sizes=tuple(chunk_sizes), 
                     overlap_frac=overlap_frac,
                     split_strategy=args.split_strategy, 
                     lowercase=lowercase,
                     batch_size=args.batch_size,
                     enable_ocr=args.enable_ocr)

if __name__ == '__main__':
    # This is necessary for multiprocessing to work correctly on
    # some platforms (like Windows)
    multiprocessing.freeze_support()
    main()