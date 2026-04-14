"""
LangChain Multi-Tools Module for Educational Multi-Agent System

Provides specialized tools for:
- Knowledge retrieval using Elasticsearch RAG
- Quiz generation using text-generation models
- Report generation using text-generation models (NEW)
- PDF document creation
- PPTX presentation generation (NEW)
- Email delivery with attachments

Each tool is decorated with @tool for seamless LangChain integration.
Models are initialized once at module load time.
All responses are JSON-formatted for consistent parsing by the main agent.
"""

import os
import re
import json
import logging
import smtplib
from pathlib import Path
from email.mime.multipart import MIMEMultipart
from email.mime.text import MIMEText
from email.mime.base import MIMEBase
from email import encoders
from langchain.tools import tool
from fpdf import FPDF
from dotenv import load_dotenv
from fpdf.enums import XPos, YPos
from unidecode import unidecode
from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN
# Gemini API imports
# REMOVED: import google.generativeai as genai
from langchain_google_genai import ChatGoogleGenerativeAI  # ADDED
from langchain_core.messages import HumanMessage           # ADDED

# Load environment variables from .env file
load_dotenv()

# NEW Import: Added for presentation generation
try:
    from pptx import Presentation
    from pptx.util import Inches, Pt
except ImportError:
    print("WARNING: `python-pptx` not installed. Presentation generator will fail.")
    print("Please run: pip install python-pptx")

# ======================== LOGGING CONFIGURATION ========================

def _setup_logger():
    """Configure logger for tools module."""
    log_dir = Path("logs")
    log_dir.mkdir(exist_ok=True)
    
    logger = logging.getLogger("multitools")
    if not logger.handlers:
        logger.setLevel(logging.INFO)
        handler = logging.FileHandler(log_dir / "tools.log", encoding="utf-8")
        formatter = logging.Formatter(
            "%(asctime)s - %(name)s - %(levelname)s - %(message)s"
        )
        handler.setFormatter(formatter)
        logger.addHandler(handler)
    
    return logger

logger = _setup_logger()

# ======================== GLOBAL INITIALIZATION ========================

# RAG Retriever - initialize once
_retriever = None

def _init_retriever():
    """Initialize RAG retriever from retrieval.py module."""
    global _retriever
    if _retriever is None:
        try:
            from retrieval import Retriever
            _retriever = Retriever(log_dir=Path("logs"))
        except Exception as e:
            logger.error(f"Failed to initialize retriever: {e}", exc_info=True)
            _retriever = None
    return _retriever

# Gemini API Client - initialize once
_gemini_chat_model = None  # CHANGED: Was _gemini_client
_gemini_model_name = None

# CHANGED: Replaced _init_gemini_client and _get_gemini_model with a single function
def _get_gemini_chat_model():
    """Initialize and return the LangChain ChatGoogleGenerativeAI model."""
    global _gemini_chat_model, _gemini_model_name
    if _gemini_chat_model is None:
        try:
            # Get API key from .env file (loaded by load_dotenv() above)
            api_key = os.getenv("GEMINI_API_KEY")
            if not api_key:
                logger.error("GEMINI_API_KEY not found in .env file!")
                return None
            
            # Get model name from .env or use default
            _gemini_model_name = os.getenv("GEMINI_MODEL", "gemini-2.5-flash")
            
            # Initialize the LangChain chat model
            _gemini_chat_model = ChatGoogleGenerativeAI(
                model=_gemini_model_name,
                google_api_key=api_key
            )
            
            logger.info(f"LangChain ChatGoogleGenerativeAI initialized with model: {_gemini_model_name}")
            
        except Exception as e:
            logger.error(f"Failed to initialize ChatGoogleGenerativeAI: {e}", exc_info=True)
            _gemini_chat_model = None # Ensure it stays None on failure
    
    return _gemini_chat_model

# ======================== TOOLS ========================

@tool
def knowledge_retrieval(query: str) -> list[str]:
    """
    Retrieve chemistry knowledge from indexed database using RAG.
    ...
    """
    try:
        if not query or len(query.strip()) < 3:
            error_msg = "Query must be at least 3 characters long."
            logger.warning(f"knowledge_retrieval: {error_msg}")
            return {"status": "error", "message": error_msg} # <-- Return full dict
        
        retriever = _init_retriever()
        if retriever is None:
            error_msg = "RAG retriever not initialized. Check Elasticsearch connection."
            logger.error(f"knowledge_retrieval: {error_msg}")
            return {"status": "error", "message": "Error: Knowledge retrieval unavailable. Check logs for details."} # <-- Return full dict
        
        results = retriever.search(query, top_k=20, rerank_top_n=5)
        
        # --- FIX: Offload models after search ---
        try:
            retriever.offload_models_to_cpu()
            logger.info("knowledge_retrieval: Offloaded retriever models to CPU.")
        except Exception as offload_e:
            logger.warning(f"knowledge_retrieval: Failed to offload models: {offload_e}")
        # --- END FIX ---
        
        if not results:
            logger.warning(f"knowledge_retrieval: No results found for query: {query}")
            return {"status": "success", "result": "No relevant information found in knowledge base."} # <-- Return full dict
        
        context_chunks = []
        for chunk in results:
            if "_source" in chunk and "chunk_text" in chunk["_source"]:
                context_chunks.append(chunk["_source"]["chunk_text"])
        
        if not context_chunks:
            error_msg = "Retrieved chunks have unexpected format."
            logger.error(f"knowledge_retrieval: {error_msg}")
            return {"status": "error", "message": "Error: Retrieved content format invalid. Check logs for details."} # <-- Return full dict
        
        logger.info(f"knowledge_retrieval: Successfully retrieved {len(context_chunks)} chunks for query: {query}")
        
        # Return a single joined string of knowledge
        joined_context = "\n\n---\n\n".join(context_chunks)
        return {"status": "success", "result": joined_context}
        
    except Exception as e:
        logger.error(f"knowledge_retrieval exception: {e}", exc_info=True)
        # --- FIX: Ensure offload even on error ---
        try:
            if _retriever:
                _retriever.offload_models_to_cpu()
                logger.info("knowledge_retrieval (on error): Offloaded retriever models to CPU.")
        except Exception as offload_e:
            logger.warning(f"knowledge_retrieval (on error): Failed to offload models: {offload_e}")
        # --- END FIX ---
        return {"status": "error", "message": "Error: Knowledge retrieval failed. Check logs for details."}

@tool
def quiz_generator(context: str, constraints: str) -> str:
    """
    Generate educational quiz content from knowledge context using Gemini API.
    
    Uses Google's Gemini API to generate quiz questions with multiple choice 
    answers based on provided context and educational constraints.
    
    Args:
        context: Knowledge content to generate quiz from
        constraints: Educational requirements (e.g., "grade 8, 10 marks, 5 questions")
        
    Returns:
        JSON string with status and generated quiz, or error message
    """
    try:
        # Validate inputs
        if not context or len(context.strip()) < 50:
            error_msg = "Context must be at least 50 characters long."
            logger.warning(f"quiz_generator: {error_msg}")
            return json.dumps({"status": "error", "message": error_msg})
        
        if not constraints or len(constraints.strip()) < 5:
            error_msg = "Constraints must be at least 5 characters long."
            logger.warning(f"quiz_generator: {error_msg}")
            return json.dumps({"status": "error", "message": error_msg})
        
        # Initialize Gemini model
        model = _get_gemini_chat_model() # CHANGED
        if model is None:
            error_msg = "Gemini API not initialized."
            logger.error(f"quiz_generator: {error_msg}")
            return json.dumps({
                "status": "error",
                "message": "Error: Quiz generation unavailable. Check GEMINI_API_KEY and logs for details."
            })
        
        # Build prompt for quiz generation
        prompt = f"""Generate an educational quiz based on the provided knowledge base.

Requirements: {constraints}

Knowledge Base:
{context[:3000]}

Create a structured quiz with:
- Clear numbered questions (at least 5)
- 4 multiple choice options (A, B, C, D) for each question
- Correct answer marked as "Answer: [letter]" after each question
- Professional formatting suitable for students
- Difficulty appropriate for the grade level specified

Begin quiz:"""
        
        logger.info(f"quiz_generator: Sending request to Gemini API with constraints: {constraints}")
        
        # Call Gemini API using LangChain interface
        messages = [HumanMessage(content=prompt)]  # ADDED
        
        # CHANGED: Switched to model.invoke()
        response = model.invoke(
            messages,
            generation_config={
                "temperature": 0.7,
                "top_p": 0.95,
                "max_output_tokens": 2048,
            }
        )
        
        quiz_content = response.content.strip() # CHANGED: from response.text
        
        # Validate output
        if not quiz_content or len(quiz_content) < 50:
            error_msg = "Quiz generation produced insufficient content."
            logger.warning(f"quiz_generator: {error_msg}")
            return json.dumps({
                "status": "error",
                "message": "Error: Generated quiz is too short. Check logs for details."
            })
        
        logger.info(f"quiz_generator: Successfully generated quiz ({len(quiz_content)} characters) for constraints: {constraints}")
        
        return json.dumps({
            "status": "success",
            "result": quiz_content,
            "length": len(quiz_content)
        })
        
    except Exception as e:
        logger.error(f"quiz_generator exception: {e}", exc_info=True)
        return json.dumps({
            "status": "error",
            "message": "Error: Quiz generation failed. Check logs for details."
        })

# --- NEW: ProfessionalPDF Class (using FPDF) ---
class ProfessionalPDF(FPDF):
    """Custom PDF class for professional headers/footers."""
    def __init__(self, title_doc):
        super().__init__()
        self.doc_title = title_doc
        self.set_left_margin(15)
        self.set_right_margin(15)
        self.set_top_margin(15)
        self.font_family = "Helvetica" 

    def header(self):
        self.set_font(self.font_family, 'B', 15)
        safe_header = unidecode('Chemistry SME Agent')
        self.cell(0, 10, safe_header, 0, 0, 'C')
        self.ln(10)
        self.set_line_width(0.5)
        self.set_draw_color(150, 150, 150)
        self.line(15, 25, self.w - 15, 25)
        self.ln(10)

    def footer(self):
        self.set_y(-15)
        self.set_font(self.font_family, 'I', 8)
        self.set_text_color(128)
        safe_footer = unidecode(f'Page {self.page_no()} | Generated by AI Educational Assistant')
        self.cell(0, 10, safe_footer, 0, 0, 'C')
# --- NEW TOOL ---
@tool
def report_generator(context: str, topic: str) -> str:
    """
    Generates a written report/summary from a knowledge context using Gemini API.
    
    Uses Google's Gemini API to synthesize the provided context
    into a coherent, multi-paragraph report on the given topic.
    
    Args:
        context: Knowledge content to generate the report from.
        topic: The central topic of the report (e.g., "The French Revolution").
        
    Returns:
        JSON string with status and the generated report, or error message.
    """
    try:
        if not context or len(context.strip()) < 50:
            error_msg = "Context must be at least 50 characters long."
            logger.warning(f"report_generator: {error_msg}")
            return json.dumps({"status": "error", "message": error_msg})
        
        if not topic or len(topic.strip()) < 3:
            error_msg = "Topic must be at least 3 characters long."
            logger.warning(f"report_generator: {error_msg}")
            return json.dumps({"status": "error", "message": error_msg})
        
        # Initialize Gemini model
        model = _get_gemini_chat_model() # CHANGED
        if model is None:
            error_msg = "Gemini API not initialized."
            logger.error(f"report_generator: {error_msg}")
            return json.dumps({
                "status": "error",
                "message": "Error: Report generation unavailable. Check GEMINI_API_KEY and logs for details."
            })
        
        # Build prompt for report generation
        prompt = f"""Generate a detailed, well-structured academic report.

Topic: {topic}

Knowledge Base:
{context[:4000]}

Write a comprehensive report based ONLY on the knowledge base provided. 

Structure your report with:
- Clear introduction paragraph
- Multiple well-developed body paragraphs covering different aspects
- Logical flow between paragraphs
- A conclusion paragraph that summarizes key points
- Professional academic tone
- Proper citations or references where possible

Begin report:"""
        
        logger.info(f"report_generator: Sending request to Gemini API for topic: {topic}")
        
        # Call Gemini API using LangChain interface
        messages = [HumanMessage(content=prompt)] # ADDED
        
        # CHANGED: Switched to model.invoke()
        response = model.invoke(
            messages,
            generation_config={
                "temperature": 0.7,
                "top_p": 0.95,
                "max_output_tokens": 3000,
            }
        )
        
        report_content = response.content.strip() # CHANGED: from response.text
        
        # Validate output
        if not report_content or len(report_content) < 100:
            error_msg = "Report generation produced insufficient content."
            logger.warning(f"report_generator: {error_msg}")
            return json.dumps({
                "status": "error",
                "message": "Error: Generated report is too short. Check logs for details."
            })
        
        logger.info(f"report_generator: Successfully generated report ({len(report_content)} characters) for topic: {topic}")
        
        return json.dumps({
            "status": "success",
            "result": report_content,
            "length": len(report_content)
        })
        
    except Exception as e:
        logger.error(f"report_generator exception: {e}", exc_info=True)
        return json.dumps({
            "status": "error",
            "message": "Error: Report generation failed. Check logs for details."
        })

@tool
def pdf_generator(content: str, filename: str, title: str):
    """
    Generate a clean, nicely formatted PDF from Markdown content
    using markdown_pdf (MarkdownPdf, Section).
    """

    try:
        # Resolve output directory (same logic as app.py)
        script_dir = os.path.dirname(os.path.abspath(__file__))
        output_dir = os.path.abspath(os.path.join(script_dir, "..", "generated_documents"))
        os.makedirs(output_dir, exist_ok=True)

        # Ensure filename ends with .pdf
        if not filename.lower().endswith(".pdf"):
            filename += ".pdf"

        filepath = os.path.join(output_dir, filename)

        # --- Clean Markdown Input ---
        content = content.replace("\\n", "\n").replace("\\t", " ").replace("\r", "")

        # OPTIONAL: Prepend a title (as an H1 heading)
        md_content = f"# {title}\n\n{content}"

        # --- Use markdown_pdf library ---
        try:
            from markdown_pdf import MarkdownPdf, Section
        except ImportError:
            logger.error("markdown_pdf is not installed. Run: pip install markdown-pdf")
            return {"status": "error", "message": "markdown_pdf not installed."}

        pdf = MarkdownPdf()

        # Add entire markdown content as one section
        pdf.add_section(Section(md_content))

        # Save PDF
        pdf.save(filepath)

        logger.info(f"✅ Markdown PDF generated successfully: {filepath}")

        return {
            "status": "success",
            "result": filepath,
            "file_size": os.path.getsize(filepath)
        }

    except Exception as e:
        logger.error(f"pdf_generator exception: {e}", exc_info=True)
        return {"status": "error", "message": "Error: PDF generation failed."}
    
# --- UPGRADED: presentation_generator (JUNK CLEANING) ---
@tool
def presentation_generator(topic: str, content: str, filename: str) -> str:
    """
    Generates a structured and STYLED .pptx presentation file from content.
    Includes logic to clean junk text before the first header.
    """
    try:
        if not filename.endswith(".pptx"):
            filename += ".pptx"

        # Use your existing path logic
        script_dir = os.path.dirname(os.path.abspath(__file__))
        output_dir = os.path.abspath(os.path.join(script_dir, "..", "generated_documents"))
        
        os.makedirs(output_dir, exist_ok=True)
        filepath = os.path.join(output_dir, filename)

        prs = Presentation()
        # Set standard slide size (16:9)
        prs.slide_width = Inches(16)
        prs.slide_height = Inches(9)
        
        # Define Colors
        COLOR_DARK_BLUE = RGBColor(0, 51, 102)
        COLOR_MEDIUM_BLUE = RGBColor(0, 102, 204)
        COLOR_TEXT_DARK = RGBColor(50, 50, 50)
        COLOR_TEXT_LIGHT = RGBColor(150, 150, 150)
        
        # --- Slide 1: Title Slide ---
        title_slide_layout = prs.slide_layouts[0] # Title slide
        slide = prs.slides.add_slide(title_slide_layout)
        
        # Title
        title_shape = slide.shapes.title
        title_shape.text = topic
        p = title_shape.text_frame.paragraphs[0]
        p.font.color.rgb = COLOR_DARK_BLUE
        p.font.bold = True
        p.font.size = Pt(54)
        
        # Subtitle
        subtitle_shape = slide.placeholders[1]
        subtitle_shape.text = "Generated by the AI Chemistry Agent"
        p = subtitle_shape.text_frame.paragraphs[0]
        p.font.color.rgb = COLOR_MEDIUM_BLUE
        p.font.size = Pt(24)

        # --- Content Parsing Logic ---
        content_layout = prs.slide_layouts[1] # Title and Content
        current_slide = None
        
        # Cleanup content
        content = content.replace("\\n", "\n").replace("\\t", " ").replace("\r", "")
        
        # --- PPTX FIX: Find first real header and slice content ---
        first_header_index = content.find("## ")
        if first_header_index != -1:
            content = content[first_header_index:]
        else:
            # If no "## " header, try to clean up known junk
            content = re.sub(r".*={5,}\n+", "", content) # Remove "=====" lines
            content = re.sub(r"^D BLOCK ELEMENTS.*?\n", "", content, flags=re.IGNORECASE)
        # --- END PPTX FIX ---
        
        lines = content.split('\n')
        
        for line in lines:
            line = line.strip()
            if not line or "===" in line: # Skip empty lines and dividers
                continue
            
            # Detect Header -> New Slide
            if line.startswith('##') or line.startswith('###') or (line.endswith(':') and len(line) < 100):
                title_text = line.replace('#', '').strip()
                if len(title_text) > 80: title_text = title_text[:77] + "..."
                
                current_slide = prs.slides.add_slide(content_layout)
                
                # Style Title
                title_shape = current_slide.shapes.title
                title_shape.text = title_text
                p = title_shape.text_frame.paragraphs[0]
                p.font.color.rgb = COLOR_DARK_BLUE
                p.font.size = Pt(36)
                p.font.bold = True
                
                # Reset text frame for body
                body_shape = current_slide.placeholders[1]
                tf = body_shape.text_frame
                tf.clear()
                tf.word_wrap = True # Ensure text wraps
                
                # --- Add Professional Footer to Content Slides ---
                txBox = current_slide.shapes.add_textbox(
                    Inches(0.5), Inches(8.5), Inches(15), Inches(0.4)
                )
                tf_footer = txBox.text_frame
                p_footer = tf_footer.paragraphs[0]
                p_footer.text = f"{topic} | Chemistry SME Agent"
                p_footer.font.size = Pt(12)
                p_footer.font.color.rgb = COLOR_TEXT_LIGHT
                p_footer.alignment = PP_ALIGN.RIGHT
                
            # Detect Bullet Points
            elif (line.startswith('- ') or line.startswith('* ')) and current_slide:
                tf = current_slide.placeholders[1].text_frame
                p = tf.add_paragraph()
                p.text = line[2:]
                p.level = 0
                p.font.size = Pt(24)
                p.font.color.rgb = COLOR_TEXT_DARK
                
            # Numbered Lists or Quiz Options
            elif (line[0].isdigit() and line[1] == '.') or line.startswith('☐') and current_slide:
                tf = current_slide.placeholders[1].text_frame
                p = tf.add_paragraph()
                p.text = line
                p.level = 0
                p.font.size = Pt(24)
                p.font.color.rgb = COLOR_TEXT_DARK
                
            # Standard text -> Add as bullet point
            elif line and current_slide: # Check for 'line' to avoid adding empty paragraphs
                if len(current_slide.placeholders[1].text_frame.text) < 1000: # Avoid overcrowding
                    tf = current_slide.placeholders[1].text_frame
                    p = tf.add_paragraph()
                    p.text = line
                    p.level = 0
                    p.font.size = Pt(22)
                    p.font.color.rgb = COLOR_TEXT_DARK
                    
            # Handle content before first header (e.g. Introduction)
            elif line and not current_slide:
                # Create an "Overview" slide
                current_slide = prs.slides.add_slide(content_layout)
                current_slide.shapes.title.text = "Overview"
                p = current_slide.shapes.title.text_frame.paragraphs[0]
                p.font.color.rgb = COLOR_DARK_BLUE
                p.font.size = Pt(36)
                p.font.bold = True
                
                tf = current_slide.placeholders[1].text_frame
                tf.clear()
                p = tf.add_paragraph()
                p.text = line
                p.font.size = Pt(24)
                p.font.color.rgb = COLOR_TEXT_DARK

                # Add footer to this first slide too
                txBox = current_slide.shapes.add_textbox(
                    Inches(0.5), Inches(8.5), Inches(15), Inches(0.4)
                )
                tf_footer = txBox.text_frame
                p_footer = tf_footer.paragraphs[0]
                p_footer.text = f"{topic} | Chemistry SME Agent"
                p_footer.font.size = Pt(12)
                p_footer.font.color.rgb = COLOR_TEXT_LIGHT
                p_footer.alignment = PP_ALIGN.RIGHT

        prs.save(filepath)
        logger.info(f"✅ Professional PPTX (v3 - Cleaned) generated: {filepath}")

        return json.dumps({"status": "success", "result": filepath, "file_size": os.path.getsize(filepath)})

    except Exception as e:
        logger.error(f"presentation_generator exception: {e}", exc_info=True)
        return json.dumps({"status": "error", "message": f"Error: Presentation generation failed."})
    
    
# --- REVERTED: email_tool (Parses Markdown) ---
@tool
def email_tool(to_email: str, subject: str, body: str, attachment_paths: list[str] | None = None) -> str:
    """
    Sends an HTML-formatted email with optional attachments.
    Renders basic Markdown (##, ###, *, ☐) into HTML.
    """
    try:
        if not to_email or not subject or not body:
            return json.dumps({"status": "error", "message": "To/Subject/Body fields are required."})
        
        smtp_server = os.getenv("EMAIL_SMTP_SERVER", "smtp.gmail.com")
        smtp_port = int(os.getenv("EMAIL_SMTP_PORT", "587"))
        sender_email = os.getenv("EMAIL_USERNAME", "periodictablesme@gmail.com")
        sender_password = os.getenv("EMAIL_PASSWORD", "rjxl wmay mmuk yxav")
        
        if not sender_email or not sender_password:
             return json.dumps({"status": "error", "message": "Email credentials missing."})
        
        msg = MIMEMultipart("alternative")
        msg["From"] = sender_email
        msg["To"] = to_email
        msg["Subject"] = subject

        # --- HTML Template ---
        formatted_body = body.replace("\n", "<br>")
        
        # --- HTML ENHANCEMENT: Render Markdown ---
        # Headers (##)
        formatted_body = re.sub(r"## (.*?)<br>", r'<h2 style="color: #003366; margin-bottom: 0; margin-top: 15px;">\1</h2>', formatted_body)
        # Sub-Headers (###)
        formatted_body = re.sub(r"### (.*?)<br>", r'<h3 style="color: #333; margin-bottom: 0; margin-top: 10px;">\1</h3>', formatted_body)
        # Quiz Options (☐)
        formatted_body = re.sub(r"(☐\s+[A-D]\).*)<br>", r'<div style="margin-left: 20px; padding: 4px; font-family: monospace; font-size: 1.1em;">\1</div>', formatted_body)
        # Bullets (* or -)
        formatted_body = re.sub(r"[\*|-]\s+(.*?)<br>", r'<div style="margin-left: 15px; padding-left: 5px;">• \1</div>', formatted_body)
        # --- END HTML ENHANCEMENT ---
        
        html_content = f"""
        <html>
          <head>
            <style>
              body {{ font-family: Arial, sans-serif; color: #333; line-height: 1.6; }}
              .container {{ width: 100%; max-width: 600px; margin: 0 auto; padding: 20px; border: 1px solid #eee; border-radius: 5px; }}
              .header {{ background-color: #003366; color: white; padding: 15px; text-align: center; border-radius: 5px 5px 0 0; }}
              .content {{ padding: 20px; background-color: #f9f9ff; }}
              .footer {{ font-size: 12px; color: #777; text-align: center; margin-top: 20px; border-top: 1px solid #ddd; padding-top: 10px; }}
            </style>
          </head>
          <body>
            <div class="container">
              <div class="header">
                <h2>Chemistry Learning Assistant</h2>
              </div>
              <div class="content">
                <p>Hello,</p>
                {formatted_body}
                <br>
                <p><strong>Please find your requested documents attached.</strong></p>
              </div>
              <div class="footer">
                <p>&copy; 2025 AI Education Systems. All rights reserved.</p>
                <p>Generated automatically by your AI Agent.</p>
              </div>
            </div>
          </body>
        </html>
        """
        
        msg.attach(MIMEText(body, "plain")) # Plain text fallback
        msg.attach(MIMEText(html_content, "html"))
        
        attachment_info = ""
        
        if attachment_paths:
            for path_str in attachment_paths:
                attachment_path = Path(path_str)
                if not attachment_path.exists():
                    logger.warning(f"email_tool: Attachment file not found, skipping: {attachment_path}")
                    continue
                
                try:
                    with open(attachment_path, "rb") as attachment:
                        part = MIMEBase("application", "octet-stream")
                        part.set_payload(attachment.read())
                    encoders.encode_base64(part)
                    part.add_header("Content-Disposition", f"attachment; filename= {attachment_path.name}")
                    msg.attach(part)
                    attachment_info += f" {attachment_path.name}"
                except Exception as e:
                    logger.error(f"email_tool attachment error: {e}", exc_info=True)
            
            if attachment_info:
                attachment_info = f" with attachments:{attachment_info}"

        # SMTP send logic
        try:
            server = smtplib.SMTP(smtp_server, smtp_port)
            server.starttls()
            server.login(sender_email, sender_password)
            server.sendmail(sender_email, to_email, msg.as_string())
            server.quit()
            
            success_msg = f"HTML Email sent to {to_email}{attachment_info}"
            logger.info(f"email_tool: {success_msg}")
            
            return json.dumps({"status": "success", "result": success_msg})
            
        except smtplib.SMTPException as e:
            logger.error(f"email_tool SMTP error: {e}", exc_info=True)
            return json.dumps({"status": "error", "message": "Error: SMTP error occurred."})
        
    except Exception as e:
        logger.error(f"email_tool exception: {e}", exc_info=True)
        return json.dumps({"status": "error", "message": "Error: Email sending failed."})
# ======================== EXPORTS ========================

__all__ = [
    "knowledge_retrieval",
    "quiz_generator",
    "pdf_generator",
    "email_tool",
    "report_generator",       # NEW
    "presentation_generator"  # NEW
]